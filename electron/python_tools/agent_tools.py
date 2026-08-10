#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ELYRA agent_tools 0.7 — productividad nativa fiable
Entrada stdin JSON: {"tool": "...", "args": {...}}
Salida: una línea JSON {"ok": bool, "result": str, ...}
"""
from __future__ import annotations

import json
import os
import sys
import traceback
from datetime import datetime
from pathlib import Path


def out(ok: bool, result: str, **extra):
    payload = {"ok": ok, "result": str(result)[:6000], **extra}
    sys.stdout.write(json.dumps(payload, ensure_ascii=False) + "\n")
    sys.stdout.flush()


def docs_root() -> Path:
    home = Path.home()
    for candidate in (home / "Documents", home / "Documentos", home):
        if candidate.exists():
            return candidate
    return home


def resolve_path(p: str) -> Path:
    if not p:
        return docs_root()
    raw = str(p).strip().strip('"').strip("'")
    # Expand ~ and env
    raw = os.path.expandvars(os.path.expanduser(raw))
    path = Path(raw)
    if path.is_absolute():
        return path
    # Common aliases
    low = raw.lower().replace("\\", "/")
    home = Path.home()
    if low.startswith("informes/") or low.startswith("informes\\"):
        return docs_root() / "Informes" / Path(raw).name if "/" not in raw[9:] and "\\" not in raw[9:] else docs_root() / raw
    if low in ("descargas", "downloads"):
        return home / "Downloads"
    if low in ("documentos", "documents"):
        return docs_root()
    if low in ("escritorio", "desktop"):
        return home / "Desktop"
    return docs_root() / path


def ensure_informes() -> Path:
    d = docs_root() / "Informes"
    d.mkdir(parents=True, exist_ok=True)
    return d


def tool_health(args: dict):
    """Comprueba dependencias Python del pipeline de archivos."""
    mods = {
        "pandas": False,
        "openpyxl": False,
        "pypdf": False,
        "docx": False,
        "pptx": False,
    }
    try:
        import pandas  # noqa: F401

        mods["pandas"] = True
    except Exception:
        pass
    try:
        import openpyxl  # noqa: F401

        mods["openpyxl"] = True
    except Exception:
        pass
    try:
        import pypdf  # noqa: F401

        mods["pypdf"] = True
    except Exception:
        pass
    try:
        import docx  # noqa: F401

        mods["docx"] = True
    except Exception:
        pass
    try:
        import pptx  # noqa: F401

        mods["pptx"] = True
    except Exception:
        pass
    missing = [k for k, v in mods.items() if not v]
    ok = len(missing) == 0
    msg = "Dependencias OK: " + ", ".join(f"{k}={'sí' if v else 'NO'}" for k, v in mods.items())
    if missing:
        msg += ". Instala: pip install -r electron/python_tools/requirements.txt"
    return out(ok, msg, modules=mods, missing=missing)


def tool_scan_folder(args: dict):
    root = resolve_path(args.get("root") or str(docs_root()))
    pattern = (args.get("pattern") or "*").lower()
    exts = args.get("extensions") or [
        ".xlsx",
        ".xls",
        ".csv",
        ".pdf",
        ".docx",
        ".pptx",
        ".txt",
        ".md",
    ]
    found = []
    if not root.exists():
        return out(False, f"No existe la carpeta {root}")
    for dirpath, dirnames, files in os.walk(root):
        dirnames[:] = [d for d in dirnames if not d.startswith(".") and d.lower() not in ("node_modules", "appdata")]
        try:
            depth = Path(dirpath).relative_to(root).parts
        except ValueError:
            continue
        if len(depth) > 5:
            continue
        for name in files:
            low = name.lower()
            if pattern != "*" and pattern not in low:
                continue
            if any(low.endswith(e) for e in exts):
                found.append(str(Path(dirpath) / name))
            if len(found) >= 50:
                break
        if len(found) >= 50:
            break
    if not found:
        return out(True, f"Sin archivos coincidentes en {root}")
    return out(True, "Archivos:\n" + "\n".join(found[:50]), count=len(found))


def tool_analyze_excel(args: dict):
    try:
        import pandas as pd
    except ImportError:
        return out(False, "Falta pandas. Ejecuta: pip install pandas openpyxl")

    path = resolve_path(args.get("path") or "")
    if not path.exists():
        # Intentar búsqueda por nombre en Documentos
        name = Path(args.get("path") or "").name
        if name:
            for root in (docs_root(), docs_root() / "Informes", Path.home() / "Downloads", Path.home() / "Desktop"):
                cand = root / name
                if cand.exists():
                    path = cand
                    break
        if not path.exists():
            return out(False, f"No existe {path}. Indica la ruta completa o copia el archivo a Documentos.")

    try:
        suffix = path.suffix.lower()
        if suffix == ".csv":
            # Encoding robusto
            for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
                try:
                    df = pd.read_csv(path, encoding=enc)
                    break
                except Exception:
                    df = None
            if df is None:
                return out(False, "No pude leer el CSV (encoding)")
        else:
            # Multi-hoja: primera hoja por defecto o sheet indicado
            sheet = args.get("sheet")
            try:
                if sheet is not None:
                    df = pd.read_excel(path, sheet_name=sheet, engine="openpyxl")
                else:
                    xl = pd.ExcelFile(path, engine="openpyxl")
                    sheet_names = xl.sheet_names
                    df = xl.parse(sheet_names[0])
                    if len(sheet_names) > 1:
                        # anotar hojas disponibles
                        pass
            except Exception:
                df = pd.read_excel(path)
    except Exception as e:
        return out(False, f"Error leyendo tabla: {e}")

    summary_lines = [
        f"Archivo: {path.name}",
        f"Ruta: {path}",
        f"Filas: {len(df)}, Columnas: {len(df.columns)}",
        f"Columnas: {', '.join(map(str, df.columns.tolist()[:40]))}",
    ]
    try:
        xl2 = pd.ExcelFile(path, engine="openpyxl") if path.suffix.lower() in (".xlsx", ".xlsm") else None
        if xl2 and len(xl2.sheet_names) > 1:
            summary_lines.append("Hojas: " + ", ".join(xl2.sheet_names[:12]))
    except Exception:
        pass

    # Nulos por columna (útil laboratorio)
    try:
        nulls = df.isna().sum()
        nulls = nulls[nulls > 0]
        if len(nulls):
            summary_lines.append(
                "Columnas con vacíos: "
                + ", ".join(f"{c}={int(n)}" for c, n in nulls.head(15).items())
            )
    except Exception:
        pass

    try:
        desc = df.describe(include="all").to_string()
        summary_lines.append("Estadísticas:\n" + desc[:2500])
    except Exception:
        pass
    summary_lines.append("Muestra (8 filas):\n" + df.head(8).to_string())

    export = args.get("export")
    export_path = None
    if export is True or str(export).lower() in ("true", "1", "yes", "sí", "si"):
        out_dir = ensure_informes()
        export_path = out_dir / f"analisis_{path.stem}_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"
        try:
            with pd.ExcelWriter(export_path, engine="openpyxl") as writer:
                df.head(1000).to_excel(writer, sheet_name="Datos", index=False)
                try:
                    df.describe().to_excel(writer, sheet_name="Resumen")
                except Exception:
                    pass
            summary_lines.append(f"Exportado: {export_path}")
        except Exception as e:
            summary_lines.append(f"No se pudo exportar: {e}")

    return out(True, "\n".join(summary_lines), path=str(export_path) if export_path else str(path))


def tool_summarize_pdf(args: dict):
    path = resolve_path(args.get("path") or "")
    if not path.exists():
        return out(False, f"No existe {path}")
    max_pages = int(args.get("max_pages") or 40)
    max_pages = max(1, min(max_pages, 80))
    text = ""
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        n_pages = len(reader.pages)
        for page in reader.pages[:max_pages]:
            try:
                text += (page.extract_text() or "") + "\n"
            except Exception:
                continue
    except Exception as e:
        return out(False, f"Error PDF (pip install pypdf): {e}")

    text = " ".join(text.split())
    if not text:
        return out(
            False,
            "PDF sin texto extraíble (puede ser escaneado/imagen). Necesitarías OCR aparte.",
        )
    chunk = text[:5000]
    return out(
        True,
        f"PDF {path.name} · páginas leídas hasta {max_pages}/{n_pages} · {len(text)} chars:\n{chunk}",
        pages=n_pages,
    )


def tool_read_docx(args: dict):
    path = resolve_path(args.get("path") or "")
    if not path.exists():
        return out(False, f"No existe {path}")
    try:
        from docx import Document

        doc = Document(str(path))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        # tablas
        table_bits = []
        for ti, table in enumerate(doc.tables[:5]):
            rows = []
            for row in table.rows[:20]:
                rows.append(" | ".join(c.text.strip() for c in row.cells))
            if rows:
                table_bits.append(f"[Tabla {ti + 1}]\n" + "\n".join(rows))
        body = "\n".join(paras)
        if table_bits:
            body += "\n\n" + "\n\n".join(table_bits)
        body = body[:8000]
        return out(True, f"Documento {path.name}:\n{body}")
    except Exception as e:
        return out(False, f"Error docx: {e}. pip install python-docx")


def tool_write_docx(args: dict):
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        return out(False, "Falta python-docx. pip install python-docx")

    title = args.get("title") or "Informe ELYRA"
    body = args.get("body") or args.get("content") or ""
    out_name = args.get("path") or f"Informes/informe_{datetime.now().strftime('%Y%m%d_%H%M')}.docx"
    path = resolve_path(out_name)
    if path.suffix.lower() != ".docx":
        path = path.with_suffix(".docx")
    path.parent.mkdir(parents=True, exist_ok=True)

    doc = Document()
    doc.add_heading(title, 0)
    doc.add_paragraph(f"Generado por ELYRA · {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    for block in str(body).split("\n"):
        block = block.strip()
        if not block:
            continue
        if block.startswith("# "):
            doc.add_heading(block[2:], level=1)
        elif block.startswith("## "):
            doc.add_heading(block[3:], level=2)
        elif block.startswith("- ") or block.startswith("• "):
            doc.add_paragraph(block.lstrip("-• "), style="List Bullet")
        else:
            p = doc.add_paragraph(block)
            for run in p.runs:
                run.font.size = Pt(11)
    doc.save(str(path))
    return out(True, f"Word creado: {path}", path=str(path))


def tool_write_pptx(args: dict):
    try:
        from pptx import Presentation
    except ImportError:
        return out(False, "Falta python-pptx. pip install python-pptx")

    title = args.get("title") or "Presentación ELYRA"
    slides_data = args.get("slides")
    if isinstance(slides_data, str):
        try:
            slides_data = json.loads(slides_data)
        except Exception:
            slides_data = [{"title": title, "bullets": [slides_data]}]
    if not slides_data:
        body = args.get("body") or "Contenido"
        slides_data = [
            {"title": title, "bullets": [b.strip() for b in body.split("\n") if b.strip()][:8]}
        ]

    out_name = args.get("path") or f"Informes/presentacion_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    path = resolve_path(out_name)
    if path.suffix.lower() != ".pptx":
        path = path.with_suffix(".pptx")
    path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = f"ELYRA · {datetime.now().strftime('%Y-%m-%d')}"
        except Exception:
            pass

    for s in slides_data[:12]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(s.get("title") or "Apartado")
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        bullets = s.get("bullets") or s.get("points") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        for i, b in enumerate(bullets[:10]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = str(b)
            p.level = 0
    prs.save(str(path))
    return out(True, f"PowerPoint creado: {path}", path=str(path))


def tool_html_dashboard(args: dict):
    title = args.get("title") or "Dashboard ELYRA"
    body = args.get("body") or args.get("html") or "<p>Sin datos</p>"
    out_name = args.get("path") or f"Informes/dashboard_{datetime.now().strftime('%Y%m%d_%H%M')}.html"
    path = resolve_path(out_name)
    if path.suffix.lower() not in (".html", ".htm"):
        path = path.with_suffix(".html")
    path.parent.mkdir(parents=True, exist_ok=True)
    safe_title = str(title).replace("<", "")
    html = f"""<!DOCTYPE html>
<html lang="es"><head><meta charset="utf-8"/><title>{safe_title}</title>
<style>
body{{font-family:Segoe UI,system-ui,sans-serif;margin:0;background:#0b1220;color:#e2e8f0}}
header{{padding:24px 32px;background:linear-gradient(135deg,#0ea5e9,#0369a1)}}
main{{padding:24px 32px;max-width:1100px;margin:0 auto}}
.card{{background:#111827;border:1px solid #1e293b;border-radius:12px;padding:20px;margin-bottom:16px}}
h1{{margin:0;font-size:1.5rem}} h2{{color:#7dd3fc;font-size:1.1rem}}
table{{width:100%;border-collapse:collapse}} td,th{{border:1px solid #334155;padding:8px;text-align:left}}
</style></head>
<body><header><h1>{safe_title}</h1><p>Generado por ELYRA · {datetime.now().isoformat(timespec='minutes')}</p></header>
<main class="card">{body}</main></body></html>"""
    path.write_text(html, encoding="utf-8")
    return out(True, f"Dashboard HTML: {path}", path=str(path))


TOOLS = {
    "health": tool_health,
    "scan_folder": tool_scan_folder,
    "analyze_excel": tool_analyze_excel,
    "summarize_pdf": tool_summarize_pdf,
    "read_docx": tool_read_docx,
    "write_docx": tool_write_docx,
    "write_pptx": tool_write_pptx,
    "html_dashboard": tool_html_dashboard,
}


def main():
    try:
        raw = sys.stdin.read()
        data = json.loads(raw)
        tool = data.get("tool")
        args = data.get("args") or {}
        fn = TOOLS.get(tool)
        if not fn:
            return out(False, f"Herramienta desconocida: {tool}")
        return fn(args)
    except Exception:
        return out(False, traceback.format_exc()[-1500:])


if __name__ == "__main__":
    main()
